#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from math import sqrt, sin, cos, pi as PI
import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.spec import autoshape_types

import pptx.shapes.connector

import svgwrite

import numpy as np

#===============================================================================

from drawml.formula import Geometry, radians

#===============================================================================

def DML(tag):
    return '{{http://schemas.openxmlformats.org/drawingml/2006/main}}{}'.format(tag)

#===============================================================================

# Internal PPT units are EMUs (English Metric Units)

EMU_PER_CM  = 360000
EMU_PER_IN  = 914400

DOTS_PER_IN = 96

EMU_PER_DOT = EMU_PER_IN/DOTS_PER_IN

def cm_coords(x, y):
#===================
    return (x/EMU_PER_CM, y/EMU_PER_CM)

def svg_coords(x, y):
#====================
    return (x/EMU_PER_DOT, y/EMU_PER_DOT)

def svg_units(emu):
#===================
    return emu/EMU_PER_DOT


def clamp(x):
    return max(0, min(x, 255))


def rgb_to_hex(r, g, b):
    return "#{0:02x}{1:02x}{2:02x}".format(clamp(r), clamp(g), clamp(b))


def ellipse_point(a, b, theta):
#==============================
    a_sin_theta = a*sin(theta)
    b_cos_theta = b*cos(theta)
    circle_radius = sqrt(a_sin_theta**2 + b_cos_theta**2)
    return (a*b_cos_theta/circle_radius, b*a_sin_theta/circle_radius)

#===============================================================================

class Transform(object):
    def __init__(self, shape, bbox=None):
        xfrm = shape.element.xfrm

        # From Section L.4.7.6 of ECMA-376 Part 1
        (Bx, By) = (svg_coords(xfrm.chOff.x, xfrm.chOff.y)
                        if xfrm.chOff is not None else
                    (0, 0))
        (Dx, Dy) = (svg_coords(xfrm.chExt.cx, xfrm.chExt.cy)
                        if xfrm.chExt is not None else
                    svg_coords(*bbox))
        (Bx_, By_) = svg_coords(xfrm.off.x, xfrm.off.y)
        (Dx_, Dy_) = svg_coords(xfrm.ext.cx, xfrm.ext.cy)

        theta = xfrm.rot*PI/180.0
        Fx = -1 if xfrm.flipH else 1
        Fy = -1 if xfrm.flipV else 1

        T_st = np.matrix([[Dx_/Dx,      0, Bx_ - (Dx_/Dx)*Bx] if Dx != 0 else [1, 0, Bx_],
                          [     0, Dy_/Dy, By_ - (Dy_/Dy)*By] if Dy != 0 else [0, 1, By_],
                          [     0,      0,                 1]])
        U = np.matrix([[1, 0, -(Bx_ + Dx_/2.0)],
                       [0, 1, -(By_ + Dy_/2.0)],
                       [0, 0,                1]])

        R = np.matrix([[cos(theta), -sin(theta), 0],
                       [sin(theta),  cos(theta), 0],
                       [0,                    0, 1]])
        Flip = np.matrix([[Fx,  0, 0],
                          [ 0, Fy, 0],
                          [ 0,  0, 1]])
        T_rf = U.I*R*Flip*U

        self._T = T_rf*T_st

    def svg_matrix(self):
        return (self._T[0, 0], self._T[1, 0],
                self._T[0, 1], self._T[1, 1],
                self._T[0, 2], self._T[1, 2])

# ===============================================================================


class Mpath(svgwrite.base.BaseElement, svgwrite.mixins.XLink):
    elementname = 'mpath'

    def __init__(self, href=None, **extra):
        super(Mpath, self).__init__(**extra)
        if href is not None:
            self.set_href(href)


# ===============================================================================


class SvgMaker(object):
    def __init__(self, slide, slide_number, slide_size, output_dir):

        self.enter_motion = None
        self.dynamic_path = None
        self.dynamic_svg = None
        self.slide_number = slide_number
        self._dwg = svgwrite.Drawing(filename=os.path.join(output_dir, 'slide{}.svg'.format(slide_number)),
                                     size=svg_coords(slide_size[0], slide_size[1]))
        self._dwg.defs.add(self._dwg.style('.non-scaling-stroke { vector-effect: non-scaling-stroke; }'))
        if self.svg_from_shapes(slide.shapes, self._dwg):
            self._dwg.save(pretty=True)

    def tostring(self):
        return self._dwg.tostring()

    def create_blink_animation(self):
        # Initial blink animation of dynamic shape
        blink_animate = self._dwg.animate(attributeName='opacity',
                                          values="1;0;1")
        blink_animate.set_timing(dur='1s',
                                 repeatCount='3')
        return blink_animate

    def create_motion_animation(self, svg_path, shape_id):
        # Initial motion animation of dynamic shape
        mpath = Mpath(href='#'+shape_id)
        entrance_motion = self._dwg.animateMotion()
        entrance_motion.add(mpath)
        entrance_motion.set_timing(dur='5s')
        entrance_motion.set_value(rotate='auto')
        self.enter_motion = entrance_motion

    def svg_from_shapes(self, shapes, svg_parent):
        dynamic_id = 'dynamic'
        static_id = 'static_'
        dynamic_svg = None
        for i, shape in enumerate(shapes):
            id_tag = dynamic_id if i == 0 else static_id+str(i)
            sub_svg = svgwrite.drawing.SVG(id=id_tag)

            if i == 0:
                dynamic_svg = sub_svg
            else:
                svg_parent.add(sub_svg)

            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                    or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or isinstance(shape, pptx.shapes.connector.Connector)):
                self.shape_to_svg(shape, sub_svg)

            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                svg_group = self._dwg.g()
                transform = Transform(shape)
                svg_group.matrix(*transform.svg_matrix())
                svg_parent.add(svg_group)
                self.svg_from_shapes(shape.shapes, svg_group)

            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                pass  # or recognise name of '#layer-id' and get layer name...

            else:
                print('"{}" {} not processed...'.format(shape.name, str(shape.shape_type)))

        if dynamic_svg:
            # if self.enter_motion:
            #     self.dynamic_path.add(self.enter_motion)
            dynamic_svg.add(self.create_blink_animation())
            svg_parent.add(dynamic_svg)

        return True

    def get_shape_color(self, shape):
        shape_color = 0
        # shape.fill.solid()
        if shape.fill.fore_color.type == 2:
            # SCHEME type
            shape_color = shape.fill.fore_color.theme_color
        elif shape.fill.fore_color.type == 1:
            # RGB type
            r, g, b = shape.fill.fore_color.rgb
            shape_color = rgb_to_hex(r, g, b)
            # shape_color = "#" + str(shape.fill.fore_color.rgb)
        return shape_color

    def get_shape_line_color(self, shape, shape_color):
        shape_line_color = shape_color
        if shape.line.color.type == 2:
            # SCHEME type
            shape_line_color = shape.line.color.theme_color
        if shape.line.color.type == 1:
            # RGB type
            r, g, b = shape.line.color.rgb
            shape_line_color = rgb_to_hex(r, g, b)
            # shape_line_color = "#" + str(shape.line.color.rgb)
        return shape_line_color

    def shape_to_svg(self, shape, svg_parent):
        geometry = Geometry(shape)
        is_dynamic = svg_parent.get_id() == 'dynamic'

        for path in geometry.path_list:
            bbox = (shape.width, shape.height) if path.w is None else (path.w, path.h)
            transform = Transform(shape, bbox)
            shape_color = self.get_shape_color(shape)
            shape_line_color = self.get_shape_line_color(shape, shape_color)
            svg_path = self._dwg.path(fill='none', stroke_width=3, class_='non-scaling-stroke') # id='sss'
            svg_path.matrix(*transform.svg_matrix())
            first_point = None
            current_point = None
            closed = False
            for c in path.getchildren():
                if   c.tag == DML('arcTo'):
                    wR = geometry.attrib_value(c, 'wR')
                    hR = geometry.attrib_value(c, 'hR')
                    stAng = radians(geometry.attrib_value(c, 'stAng'))
                    swAng = radians(geometry.attrib_value(c, 'swAng'))
                    p1 = ellipse_point(wR, hR, stAng)
                    p2 = ellipse_point(wR, hR, stAng + swAng)
                    pt = (current_point[0] - p1[0] + p2[0],
                          current_point[1] - p1[1] + p2[1])
                    large_arc_flag = 1 if swAng >= PI else 0
                    svg_path.push('A', svg_units(wR), svg_units(hR),
                                       0, large_arc_flag, 1,
                                       svg_units(pt[0]), svg_units(pt[1]))
                    current_point = pt

                elif c.tag == DML('close'):
                    if first_point is not None and first_point == current_point:
                        closed = True
                    svg_path.push('Z')
                    first_point = None
                elif c.tag == DML('cubicBezTo'):
                    coords = []
                    for p in c.getchildren():
                        pt = geometry.point(p)
                        coords.append(svg_units(pt[0]))
                        coords.append(svg_units(pt[1]))
                        current_point = pt
                    svg_path.push('C', *coords)
                elif c.tag == DML('lnTo'):
                    pt = geometry.point(c.pt)
                    svg_path.push('L', svg_units(pt[0]), svg_units(pt[1]))
                    current_point = pt
                elif c.tag == DML('moveTo'):
                    pt = geometry.point(c.pt)
                    svg_path.push('M', svg_units(pt[0]), svg_units(pt[1]))
                    if first_point is None:
                        first_point = pt
                    current_point = pt
                elif c.tag == DML('quadBezTo'):
                    coords = []
                    for p in c.getchildren():
                        pt = geometry.point(p)
                        coords.append(svg_units(pt[0]))
                        coords.append(svg_units(pt[1]))
                        current_point = pt
                    svg_path.push('Q', *coords)
                else:
                    print('Unknown path element: {}'.format(c.tag))

            svg_path.attribs['fill'] = shape_color
            svg_path.attribs['stroke'] = shape_line_color

            if is_dynamic:
                self.dynamic_path = svg_path
                self.dynamic_svg = svg_parent

            # Add initial motion animation if arrow is in slide (Not done yet)
            # elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            #     freeform_id = svg_parent.get_id()
            #     svg_path.attribs['id'] = freeform_id
            #     self.create_motion_animation(svg_path, freeform_id)
                # self.dynamic_svg.add(svg_path)
                # continue

            svg_parent.add(svg_path)

# ============================================================================


class SvgExtract(object):
    def __init__(self, powerpoint, output_dir):
        # self._args = args
        self._ppt = Presentation(powerpoint)
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        self._output_dir = output_dir
        self._slides = self._ppt.slides
        self._slide_size = [self._ppt.slide_width, self._ppt.slide_height]

    def slide_to_svg(self, slide_number):
        svg_maker = SvgMaker(self._slides[slide_number-1], slide_number,
                             self._slide_size, self._output_dir)

    def slides_to_svg(self):
        for n in range(1, len(self._slides)+1):
            self.slide_to_svg(n)

#===============================================================================
#
# if __name__ == '__main__':
#     import argparse
#
#     parser = argparse.ArgumentParser(description='Extract geometries from Powerpoint slides.')
#     parser.add_argument('--version', action='version', version='0.2.1')
#     parser.add_argument('--debug-xml', action='store_true',
#                         help="save a slide's DrawML for debugging")
#     parser.add_argument('--slide', type=int, metavar='N',
#                         help='only process this slide number (1-origin)')
#     parser.add_argument('output_dir', metavar='OUTPUT_DIRECTORY',
#                         help='directory in which to save geometries')
#     parser.add_argument('powerpoint', metavar='POWERPOINT_FILE',
#                         help='the name of a Powerpoint file')
#
#
#     args = parser.parse_args()
#
#     if not os.path.exists(args.output_dir):
#         os.makedirs(args.output_dir)
#
#     svg_extract = SvgExtract(args)
#     if args.slide is None:
#         svg_extract.slides_to_svg()
#     else:
#         svg_extract.slide_to_svg(args.slide)

#===============================================================================